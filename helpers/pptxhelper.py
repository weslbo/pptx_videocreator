import requests
import os
from pptx import Presentation
from pptx.util import Inches, Cm


def remove_notes(pptx):
    presentation = Presentation(pptx)
    
    for slide in presentation.slides:
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.clear()
            
        print(f"- Removing notes from slide {presentation.slides.index(slide) + 1}")
        
        presentation.save(pptx.replace(".pptx", "_no_notes.pptx"))


def addvideo1(slide, mp4):
    print("- Adding video to slide")

    left = Cm(20.28)
    top = Cm(11.41)
    width = Cm(13.59)
    height = Cm(7.64)
    
    movie = slide.shapes.add_movie(mp4, left, top, width, height, poster_frame_image=None, mime_type='video/mp4')
    movie.media_format.auto_play = True

def addvideo(slide, mp4):
    print("- Adding video to slide")

    left = Cm(20.28)
    top = Cm(11.41)
    width = Cm(13.59)
    height = Cm(7.64)
    
    movie = slide.shapes.add_movie(mp4, left, top, width, height, poster_frame_image=None, mime_type='video/mp4')
    slide.shapes._spTree.remove(movie._element)
    slide.shapes._spTree.insert(2, movie._element)
    movie.media_format.auto_play = True
    

    
