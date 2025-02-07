import sys
sys.path.append('helpers')

import pptxhelper, videohelper, gpthelper, downloadhelper
from pptx import Presentation
import re
from dotenv import load_dotenv

load_dotenv()

def is_url(text):
    url_pattern = re.compile(r'http[s]?://(?:[a-zA-Z]|[0-9]|[$-_@.&+]|[!*\\(\\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+')
    return re.match(url_pattern, text) is not None


def main():
    pptx_input_file = "pptx/DP600-01-06.pptx"
    
    presentation = Presentation(pptx_input_file)
    for slide in presentation.slides:
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text
            
            if is_url(notes_text):
                url = notes_text.strip()
                print(f"Slide {presentation.slides.index(slide) + 1}")
                
                markdown = downloadhelper.retrieve_markdown(url)
            else:
                markdown = notes_text.strip()
                       
            background_color = slide.background.fill.fore_color.rgb if slide.background.fill.type == 1 else None
            if background_color:
                hex_color = f"#{background_color}FF"
                print(f"Background color of slide {presentation.slides.index(slide) + 1}: {hex_color}")
            else:
                print(f"Slide {presentation.slides.index(slide) + 1} has no solid background color.")
                hex_color = "#FFFFFFFF"
                
            transcript = gpthelper.generate_transcript(markdown)
            ssml = gpthelper.generate_ssml(transcript)
            
            with open("./temp/ssml.xml", "w") as ssml_file:
                ssml_file.write(ssml)
            
            notes_slide.notes_text_frame.text = transcript
            video = videohelper.generate_video(ssml, hex_color)
            pptxhelper.addvideo(slide, video)

    presentation.save(pptx_input_file.replace(".pptx", "_video.pptx"))        
                
if __name__ == "__main__":
    main()